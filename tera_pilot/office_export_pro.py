"""Tera Pilot Pro feature — ``office_export_pro``.

Batch / templated export of an agent's structured result into one or more
office formats (``.docx`` / ``.xlsx`` / ``.pptx``) with a branded header /
footer, in an isolated workspace-respecting way.

Design (matches the monetization constraint "don't cut the core"):
  - The free coding core is untouched. This is a *side-channel* helper an
    agent (or the Web UI) calls explicitly to render a report; it is not
    wired into the agent loop.
  - Pro-gated via ``licensing.is_feature_licensed("office_export_pro")``,
    fail-closed: an unlicensed caller gets a structured ``{ok: False,
    error: "pro_required"}`` and ``LicenseRequiredError`` is raised when
    the caller uses the strict API. No crash, no partial files.
  - Fully offline, no telemetry, deterministic (no LLM, no network).
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional

from .licensing import LicenseRequiredError, is_feature_licensed as _licensed

#: The Pro feature id this module gates on.
FEATURE_ID = "office_export_pro"


class OfficeExportError(Exception):
    """Raised when an export bundle cannot be produced.


    Subclassed from ``Exception`` (not ``LicenseRequiredError``) so callers
    can distinguish a subscription problem from a file/template problem.
    """


@dataclass
class ReportSection:
    """One section of a structured report to render into office files."""

    title: str
    body: str = ""
    rows: Optional[List[List[Any]]] = None  # for xlsx / tables
    headers: Optional[List[str]] = None      # for xlsx header row


@dataclass
class ExportBundle:
    """The result of a successful export: a list of produced files."""

    files: List[Path] = field(default_factory=list)


def _require_license() -> None:
    """Raises ``LicenseRequiredError`` when the Pro feature isn't licensed."""
    if not _licensed(FEATURE_ID):
        raise LicenseRequiredError(
            f"{FEATURE_ID} is a Pro feature — activate a license with: "
            "tera-pilot license activate <key>"
        )


class OfficeExportPro:
    """Render a structured result into .docx / .xlsx / .pptx files.

    Stateless like ``OfficeWorker``; every call opens/saves its own files.

    Args:
        resolve_path_fn: a ``(path: str) -> Path`` callable used to keep
            every output inside the workspace (e.g. ToolEngine._resolve_path).
            Defaults to plain (absolute) Path resolution.
    """

    def __init__(self, resolve_path_fn=None):
        self._resolve = resolve_path_fn or (lambda p: Path(p).resolve())

    # ── Public API (safe) ─────────────────────────────────────────

    def export_bundle(
        self,
        *,
        title: str,
        sections: List[ReportSection],
        out_dir: str,
        formats: Optional[Dict[str, Any]] = None,
        brand: Optional[Dict[str, str]] = None,
    ) -> Dict[str, Any]:
        """Export *sections* into the requested office formats.

        Returns a structured dict (never raises for control flow):
          - ``{ok: True, files: [str, ...]}`` on success,
          - ``{ok: False, error: "pro_required"}`` when unlicensed,
          - ``{ok: False, error: <message>}`` on any other failure.

        ``formats`` selects files to write, e.g. ``{"docx": {}, "xlsx": {}}``
        (empty dict = defaults). ``brand`` may carry ``{"header", "footer"}``
        short lines inserted into the .docx.
        """
        try:
            _require_license()
        except LicenseRequiredError:
            return {"ok": False, "error": "pro_required"}
        try:
            out_root = self._resolve(out_dir)
            out_root.mkdir(parents=True, exist_ok=True)
            formats = formats if formats is not None else {"docx": {}, "xlsx": {}, "pptx": {}}
            brand = brand or {}
            produced: List[Path] = []
            if "xlsx" in formats:
                produced.append(self._render_xlsx(out_root, title, sections))
            if "docx" in formats:
                produced.append(self._render_docx(out_root, title, sections, brand))
            if "pptx" in formats:
                produced.append(self._render_pptx(out_root, title, sections))
            return {"ok": True, "files": [str(p) for p in produced]}
        except OfficeExportError as e:
            return {"ok": False, "error": str(e)}
        except Exception as e:  # pragma: no cover - defensive
            return {"ok": False, "error": f"export failed: {e}"}

    # ── Strict API (raises) ───────────────────────────────────────

    def require(self) -> None:
        """Explicit Pro gate. Raises ``LicenseRequiredError`` when unlicensed.

        Useful for callers that want an exception (tests, API layer).
        """
        _require_license()

    # ── Per-format renderers ──────────────────────────────────────

    def _slug(self, text: str) -> str:
        s = re.sub(r"[^A-Za-z0-9_\-]+", "_", text.strip().lower())
        return s.strip("_") or "report"

    def _render_xlsx(self, out_root: Path, title: str,
                     sections: List[ReportSection]) -> Path:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill

        stem = self._slug(title)
        p = out_root / f"{stem}.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = self._slug(title)[:31] or "Report"
        ws.append([title])
        ws["A1"].font = Font(bold=True, size=14)
        sheet_header = PatternFill("solid", fgColor="1F4E79")
        for sec in sections:
            ws.append([])
            ws.append([sec.title])
            ws.cell(row=ws.max_row, column=1).font = Font(bold=True)
            if sec.headers:
                ws.append(list(sec.headers))
                for cell in ws[ws.max_row]:
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.fill = sheet_header
            if sec.rows:
                for row in sec.rows:
                    ws.append(list(row) if isinstance(row, (list, tuple)) else [row])
            elif sec.body:
                ws.append([sec.body])
        ws.freeze_panes = "A2"
        wb.save(str(p))
        return p

    def _render_docx(self, out_root: Path, title: str,
                     sections: List[ReportSection],
                     brand: Dict[str, str]) -> Path:
        from docx import Document
        from docx.shared import Pt

        stem = self._slug(title)
        p = out_root / f"{stem}.docx"
        doc = Document()
        header = doc.sections[0].header
        if brand.get("header"):
            header.paragraphs[0].text = brand["header"]
        if brand.get("footer"):
            doc.sections[0].footer.paragraphs[0].text = brand["footer"]
        doc.add_heading(title, level=0)
        for sec in sections:
            doc.add_heading(sec.title, level=1)
            if sec.body:
                doc.add_paragraph(sec.body)
            if sec.headers and sec.rows:
                tbl = doc.add_table(rows=1, cols=len(sec.headers))
                hdr = tbl.rows[0].cells
                for i, h in enumerate(sec.headers):
                    hdr[i].text = str(h)
                for row in sec.rows:
                    cells = tbl.add_row().cells
                    for i, val in enumerate(
                        list(row) if isinstance(row, (list, tuple)) else [row]
                    ):
                        if i < len(cells):
                            cells[i].text = str(val)
        doc.save(str(p))
        return p

    def _render_pptx(self, out_root: Path, title: str,
                     sections: List[ReportSection]) -> Path:
        from pptx import Presentation
        from pptx.util import Pt

        stem = self._slug(title)
        p = out_root / f"{stem}.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        for sec in sections:
            s = prs.slides.add_slide(prs.slide_layouts[1])
            s.shapes.title.text = sec.title
            tf = s.placeholders[1].text_frame
            if sec.body:
                tf.text = sec.body
            if sec.rows:
                for row in sec.rows:
                    tf.add_paragraph().text = " | ".join(
                        str(v) for v in (list(row) if isinstance(row, (list, tuple)) else [row])
                    )
        prs.save(str(p))
        return p