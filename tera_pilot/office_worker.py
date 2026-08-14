"""
Tera Pilot v1.2.0 — Office Worker module.

A self-contained Office document engine for the `office` section. Built
from scratch for Tera Pilot — it does NOT shell out to any external CLI and
does NOT embed third-party agent code. The conceptual surface (create,
view, add paragraph / heading / table / slide / shape / chart, set cell,
find/replace, save_as) is intentionally similar in spirit to common
Office CLIs, but the implementation is a fresh Python wrapper over three
well-known libraries:

  - python-docx   → .docx
  - openpyxl      → .xlsx
  - python-pptx   → .pptx

Design rules:
  1. Every public method returns a short, human-readable status string
     suitable for an LLM observation ("[CREATED] report.docx",
     "[ADDED PARAGRAPH] report.docx (style=Heading 1)"). The agent
     consumes these as tool results.
  2. Errors are caught and returned as "[OFFICE ERROR] ..." strings —
     never raised. The agent can react and retry.
  3. Files are addressed by absolute or workspace-relative path. All
     writes go through the parent ToolEngine's `_resolve_path()` so the
     existing workspace sandbox applies — no path escapes the project.
  4. Stateful documents are NOT held in memory between calls. Each
     operation opens the file, mutates it, and saves it back. This
     sacrifices some performance but makes the tool stateless and
     crash-resilient — exactly what an agent loop needs.
  5. Only the operations an agent is likely to need are exposed. There
     is no attempt to mirror the full OpenXML surface; raw XML editing
     is left to the user's own code via the existing run_code tool.
"""

from __future__ import annotations

import logging
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)


# ── Format detection ───────────────────────────────────────────────────

def _detect_format(path: str) -> str:
    """Return 'docx' | 'xlsx' | 'pptx' based on file extension."""
    ext = Path(path).suffix.lower().lstrip(".")
    if ext in ("docx", "word"):
        return "docx"
    if ext in ("xlsx", "xls", "excel"):
        return "xlsx"
    if ext in ("pptx", "ppt", "powerpoint"):
        return "pptx"
    raise ValueError(
        f"[OFFICE ERROR] unsupported file extension: {ext!r} "
        f"(expected .docx / .xlsx / .pptx)"
    )


def _require_lib(format_id: str) -> None:
    """Lazy-import the library for the given format and raise a friendly
    error if it isn't installed."""
    if format_id == "docx":
        try:
            import docx  # noqa: F401
        except ImportError as e:
            raise RuntimeError(
                "python-docx is not installed. Run `pip install python-docx`."
            ) from e
    elif format_id == "xlsx":
        try:
            import openpyxl  # noqa: F401
        except ImportError as e:
            raise RuntimeError(
                "openpyxl is not installed. Run `pip install openpyxl`."
            ) from e
    elif format_id == "pptx":
        try:
            import pptx  # noqa: F401
        except ImportError as e:
            raise RuntimeError(
                "python-pptx is not installed. Run `pip install python-pptx`."
            ) from e


# ── Color helpers ──────────────────────────────────────────────────────

_NAMED_COLORS: Dict[str, str] = {
    "black": "000000", "white": "FFFFFF", "red": "FF0000",
    "green": "00FF00", "blue": "0000FF", "yellow": "FFFF00",
    "cyan": "00FFFF", "magenta": "FF00FF", "gray": "808080",
    "grey": "808080", "darkred": "8B0000", "darkblue": "00008B",
    "darkgreen": "006400", "orange": "FFA500", "purple": "800080",
    "brown": "A52A2A", "pink": "FFC0CB", "lightgray": "D3D3D3",
    "lightgrey": "D3D3D3", "navy": "000080", "teal": "008080",
}


def _parse_color(value: Optional[str]) -> Optional[str]:
    """Normalize a color spec to a 6-digit uppercase RGB hex string.

    Accepts: '#FF0000', 'FF0000', 'red', 'rgb(255, 0, 0)'.
    Returns None for empty input.
    """
    if not value:
        return None
    v = value.strip().strip("#").upper()
    if v.lower() in _NAMED_COLORS:
        return _NAMED_COLORS[v.lower()]
    if v.startswith("RGB(") and v.endswith(")"):
        parts = v[4:-1].split(",")
        if len(parts) == 3:
            try:
                return "{:02X}{:02X}{:02X}".format(
                    *[int(p.strip()) for p in parts]
                )
            except ValueError:
                pass
    if re.fullmatch(r"[0-9A-Fa-f]{6}", v):
        return v.upper()
    if re.fullmatch(r"[0-9A-Fa-f]{3}", v):
        # 'F00' → 'FF0000'
        return "".join(c * 2 for c in v).upper()
    raise ValueError(f"unrecognized color: {value!r}")


# ── Unit helpers ───────────────────────────────────────────────────────

def _parse_emu(value: Any) -> Optional[int]:
    """Parse a dimension into EMU (English Metric Units), which is what
    python-pptx uses internally. Accepts:

      - int / float (assumed already in EMU)
      - '12pt', '12pt', '0.5cm', '1in', '96px'

    Returns None for empty input.
    """
    if value is None or value == "":
        return None
    if isinstance(value, (int, float)):
        return int(value)
    s = str(value).strip().lower()
    m = re.fullmatch(r"(-?\d+(?:\.\d+)?)\s*([a-z]+)", s)
    if not m:
        try:
            return int(float(s))
        except ValueError:
            raise ValueError(f"unrecognized dimension: {value!r}")
    num = float(m.group(1))
    unit = m.group(2)
    # 1 inch = 914400 EMU; 1 cm = 360000 EMU; 1 pt = 12700 EMU; 1 px = 9525 EMU
    if unit in ("in", "inch"):
        return int(num * 914400)
    if unit == "cm":
        return int(num * 360000)
    if unit == "mm":
        return int(num * 36000)
    if unit in ("pt", "pts"):
        return int(num * 12700)
    if unit in ("px", "pxl", "pixels"):
        return int(num * 9525)
    if unit == "emu":
        return int(num)
    raise ValueError(f"unrecognized unit in {value!r}")


# ── Main OfficeWorker class ────────────────────────────────────────────

class OfficeWorker:
    """Stateless office document operations for the `office` agent section.

    The class is instantiated once per ToolEngine and shared across all
    office_* tool calls. It holds no per-document state — each method
    opens, mutates, and saves the file before returning.
    """

    def __init__(self, resolve_path_fn=None):
        # ``resolve_path_fn`` is ToolEngine._resolve_path — used to keep
        # every file access inside the workspace sandbox. If not provided
        # (e.g. in tests), we fall back to Path resolution.
        self._resolve = resolve_path_fn or (lambda p: Path(p).resolve())

    # ── Create ─────────────────────────────────────────────────────

    def create(self, path: str, template: str = "blank") -> str:
        """Create a blank .docx / .xlsx / .pptx file."""
        try:
            fmt = _detect_format(path)
            _require_lib(fmt)
            p = self._resolve(path)
            p.parent.mkdir(parents=True, exist_ok=True)
            if p.exists():
                return f"[OFFICE ERROR] file already exists: {path}"
            if fmt == "docx":
                from docx import Document
                doc = Document()
                doc.save(str(p))
            elif fmt == "xlsx":
                from openpyxl import Workbook
                wb = Workbook()
                # openpyxl's default sheet is named "Sheet" — rename to
                # "Sheet1" so it matches what users see in Excel/LibreOffice
                # and so the agent can address it predictably.
                if wb.sheetnames:
                    wb[wb.sheetnames[0]].title = "Sheet1"
                wb.save(str(p))
            elif fmt == "pptx":
                from pptx import Presentation
                prs = Presentation()
                prs.save(str(p))
            return f"[CREATED] {path} ({fmt})"
        except Exception as e:
            return f"[OFFICE ERROR] create failed: {e}"

    # ── View ───────────────────────────────────────────────────────

    def view(self, path: str, mode: str = "outline") -> str:
        """Inspect a document. Modes:
          - outline: structural tree (sections, paragraphs, sheets, slides)
          - text:    plain text extraction
          - stats:   word/sheet/slide counts
        """
        try:
            fmt = _detect_format(path)
            _require_lib(fmt)
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            mode = (mode or "outline").lower()
            if fmt == "docx":
                return self._view_docx(p, mode)
            if fmt == "xlsx":
                return self._view_xlsx(p, mode)
            return self._view_pptx(p, mode)
        except Exception as e:
            return f"[OFFICE ERROR] view failed: {e}"

    def _view_docx(self, p: Path, mode: str) -> str:
        from docx import Document
        doc = Document(str(p))
        if mode == "text":
            return "\n".join(para.text for para in doc.paragraphs)
        if mode == "stats":
            n_para = len(doc.paragraphs)
            n_tables = len(doc.tables)
            n_words = sum(len(para.text.split()) for para in doc.paragraphs)
            return (
                f"DOCX stats: paragraphs={n_para}, tables={n_tables}, "
                f"words={n_words}"
            )
        # outline (default)
        lines = [f"DOCX outline ({len(doc.paragraphs)} paragraphs, "
                 f"{len(doc.tables)} tables):"]
        for i, para in enumerate(doc.paragraphs, 1):
            style = para.style.name if para.style else "None"
            text_preview = para.text[:80].replace("\n", " ")
            lines.append(f"  [{i}] ({style}) {text_preview}")
        for i, tbl in enumerate(doc.tables, 1):
            lines.append(f"  table#{i}: {len(tbl.rows)} rows x {len(tbl.columns)} cols")
        return "\n".join(lines)

    def _view_xlsx(self, p: Path, mode: str) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(str(p), read_only=True, data_only=False)
        if mode == "text":
            out = []
            for ws in wb.worksheets:
                out.append(f"## {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    out.append("\t".join("" if v is None else str(v) for v in row))
            wb.close()
            return "\n".join(out)
        if mode == "stats":
            total = sum(ws.max_row * ws.max_column for ws in wb.worksheets)
            return (
                f"XLSX stats: sheets={len(wb.sheetnames)}, "
                f"names={wb.sheetnames}, total_cells={total}"
            )
        # outline
        lines = [f"XLSX outline ({len(wb.sheetnames)} sheets):"]
        for ws in wb.worksheets:
            lines.append(f"  - {ws.title}: {ws.max_row} rows x {ws.max_column} cols")
        wb.close()
        return "\n".join(lines)

    def _view_pptx(self, p: Path, mode: str) -> str:
        from pptx import Presentation
        prs = Presentation(str(p))
        if mode == "text":
            out = []
            for i, slide in enumerate(prs.slides, 1):
                out.append(f"## Slide {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        out.append(shape.text_frame.text)
            return "\n".join(out)
        if mode == "stats":
            n_shapes = sum(len(s.shapes) for s in prs.slides)
            return (
                f"PPTX stats: slides={len(prs.slides)}, "
                f"total_shapes={n_shapes}"
            )
        # outline
        lines = [f"PPTX outline ({len(prs.slides)} slides):"]
        for i, slide in enumerate(prs.slides, 1):
            shape_summary = []
            for shape in slide.shapes:
                kind = shape.shape_type
                text = ""
                if shape.has_text_frame:
                    text = shape.text_frame.text[:40].replace("\n", " ")
                shape_summary.append(f"{kind}" + (f": {text!r}" if text else ""))
            lines.append(f"  slide {i}: {' | '.join(shape_summary) if shape_summary else '(empty)'}")
        return "\n".join(lines)

    # ── DOCX: paragraphs, headings, tables ─────────────────────────

    def add_paragraph(self, path: str, text: str, style: str = "Normal",
                      bold: bool = False, italic: bool = False,
                      color: Optional[str] = None,
                      size: Optional[int] = None) -> str:
        """Append a paragraph to a .docx. Style can be 'Normal',
        'Heading 1', 'Heading 2', 'List Bullet', etc."""
        try:
            _require_lib("docx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from docx import Document
            from docx.shared import Pt, RGBColor
            doc = Document(str(p))
            try:
                para = doc.add_paragraph(text, style=style)
            except KeyError:
                # Style doesn't exist in the document's styles — fall back
                # to Normal and apply formatting directly.
                para = doc.add_paragraph(text)
                para.style = doc.styles["Normal"]
            run = para.add_run() if not para.runs else para.runs[0]
            # Re-create the run with the text so formatting applies cleanly.
            if not para.runs:
                run = para.add_run(text)
            else:
                run = para.runs[0]
                run.text = text
            if bold:
                run.bold = True
            if italic:
                run.italic = True
            if size:
                run.font.size = Pt(int(size))
            if color:
                hexc = _parse_color(color)
                if hexc:
                    run.font.color.rgb = RGBColor.from_string(hexc)
            doc.save(str(p))
            return f"[ADDED PARAGRAPH] {path} (style={style})"
        except Exception as e:
            return f"[OFFICE ERROR] add_paragraph failed: {e}"

    def add_heading(self, path: str, text: str, level: int = 1) -> str:
        """Add a heading (level 1-9) to a .docx."""
        try:
            _require_lib("docx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from docx import Document
            doc = Document(str(p))
            level = max(1, min(int(level), 9))
            doc.add_heading(text, level=level)
            doc.save(str(p))
            return f"[ADDED HEADING] {path} (level={level})"
        except Exception as e:
            return f"[OFFICE ERROR] add_heading failed: {e}"

    def add_table(self, path: str, rows: int, cols: int,
                  data: Optional[List[List[Any]]] = None,
                  header: bool = True) -> str:
        """Add a table to a .docx. If `data` is provided, it fills the
        cells row-by-row (data[i][j]). Missing cells stay empty."""
        try:
            _require_lib("docx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from docx import Document
            doc = Document(str(p))
            rows = max(1, int(rows))
            cols = max(1, int(cols))
            tbl = doc.add_table(rows=rows, cols=cols)
            tbl.style = "Table Grid"
            if data:
                for i, row_data in enumerate(data[:rows]):
                    for j, val in enumerate(row_data[:cols]):
                        cell = tbl.cell(i, j)
                        cell.text = str(val) if val is not None else ""
            if header and data:
                # Bold the first row.
                for j in range(min(cols, len(data[0]) if data else 0)):
                    for para in tbl.cell(0, j).paragraphs:
                        for run in para.runs:
                            run.bold = True
            doc.save(str(p))
            return f"[ADDED TABLE] {path} ({rows}x{cols})"
        except Exception as e:
            return f"[OFFICE ERROR] add_table failed: {e}"

    # ── XLSX: sheets, cells, formulas ──────────────────────────────

    def add_sheet(self, path: str, name: str) -> str:
        """Add a new sheet to an .xlsx workbook."""
        try:
            _require_lib("xlsx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from openpyxl import load_workbook
            wb = load_workbook(str(p))
            if name in wb.sheetnames:
                return f"[OFFICE ERROR] sheet already exists: {name}"
            wb.create_sheet(title=name)
            wb.save(str(p))
            wb.close()
            return f"[ADDED SHEET] {path} / {name}"
        except Exception as e:
            return f"[OFFICE ERROR] add_sheet failed: {e}"

    def set_cell(self, path: str, sheet: str, cell: str, value: Any) -> str:
        """Set the value of a single cell. `value` may be a number,
        string, or a formula beginning with '='."""
        try:
            _require_lib("xlsx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from openpyxl import load_workbook
            wb = load_workbook(str(p))
            if sheet not in wb.sheetnames:
                wb.close()
                return f"[OFFICE ERROR] sheet not found: {sheet}"
            ws = wb[sheet]
            ws[cell] = value
            wb.save(str(p))
            wb.close()
            return f"[SET CELL] {path} / {sheet}!{cell} = {value!r}"
        except Exception as e:
            return f"[OFFICE ERROR] set_cell failed: {e}"

    def set_cell_format(self, path: str, sheet: str, cell: str,
                        bold: Optional[bool] = None,
                        italic: Optional[bool] = None,
                        font_color: Optional[str] = None,
                        bg_color: Optional[str] = None,
                        font_size: Optional[int] = None,
                        align: Optional[str] = None) -> str:
        """Apply formatting to a single cell."""
        try:
            _require_lib("xlsx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from openpyxl import load_workbook
            from openpyxl.styles import Font, PatternFill, Alignment
            wb = load_workbook(str(p))
            if sheet not in wb.sheetnames:
                wb.close()
                return f"[OFFICE ERROR] sheet not found: {sheet}"
            ws = wb[sheet]
            cell_obj = ws[cell]
            font_kwargs = {}
            if bold is not None:
                font_kwargs["bold"] = bool(bold)
            if italic is not None:
                font_kwargs["italic"] = bool(italic)
            if font_color:
                hexc = _parse_color(font_color)
                if hexc:
                    font_kwargs["color"] = hexc
            if font_size:
                font_kwargs["size"] = int(font_size)
            if font_kwargs:
                # Preserve existing font props not specified.
                existing = cell_obj.font
                merged = {
                    "name": existing.name, "size": existing.size,
                    "bold": existing.bold, "italic": existing.italic,
                    "color": existing.color,
                }
                merged.update(font_kwargs)
                cell_obj.font = Font(**merged)
            if bg_color:
                hexc = _parse_color(bg_color)
                if hexc:
                    cell_obj.fill = PatternFill(
                        start_color=hexc, end_color=hexc,
                        fill_type="solid",
                    )
            if align:
                align = align.lower()
                horiz = {"left": "left", "center": "center", "centre": "center",
                         "right": "right"}.get(align, None)
                if horiz:
                    cell_obj.alignment = Alignment(horizontal=horiz)
            wb.save(str(p))
            wb.close()
            return f"[SET CELL FORMAT] {path} / {sheet}!{cell}"
        except Exception as e:
            return f"[OFFICE ERROR] set_cell_format failed: {e}"

    def add_chart(self, path: str, sheet: str, chart_type: str = "bar",
                  data_range: str = "A1:B10", anchor: str = "D2",
                  title: Optional[str] = None) -> str:
        """Add a chart to an .xlsx sheet. Supported types: bar, line,
        pie, scatter."""
        try:
            _require_lib("xlsx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from openpyxl import load_workbook
            from openpyxl.chart import BarChart, LineChart, PieChart, ScatterChart, Reference
            wb = load_workbook(str(p))
            if sheet not in wb.sheetnames:
                wb.close()
                return f"[OFFICE ERROR] sheet not found: {sheet}"
            ws = wb[sheet]
            ct = (chart_type or "bar").lower()
            if ct == "bar":
                chart = BarChart()
            elif ct == "line":
                chart = LineChart()
            elif ct == "pie":
                chart = PieChart()
            elif ct == "scatter":
                chart = ScatterChart()
            else:
                wb.close()
                return f"[OFFICE ERROR] unknown chart_type: {chart_type}"
            # Parse data_range like "A1:B10"
            try:
                start, end = data_range.split(":")
            except ValueError:
                wb.close()
                return f"[OFFICE ERROR] data_range must be like 'A1:B10'"
            from openpyxl.utils import range_boundaries
            min_col, min_row, max_col, max_row = range_boundaries(
                f"{start}:{end}"
            )
            data_ref = Reference(
                ws, min_col=min_col, min_row=min_row,
                max_col=max_col, max_row=max_row,
            )
            if ct == "scatter":
                chart.xvalues = data_ref
            else:
                chart.add_data(data_ref, titles_from_data=False)
            if title:
                chart.title = title
            ws.add_chart(chart, anchor)
            wb.save(str(p))
            wb.close()
            return f"[ADDED CHART] {path} / {sheet} ({ct} @ {anchor})"
        except Exception as e:
            return f"[OFFICE ERROR] add_chart failed: {e}"

    # ── PPTX: slides, text, shapes ─────────────────────────────────

    def add_slide(self, path: str, layout: str = "title",
                  title: Optional[str] = None,
                  subtitle: Optional[str] = None) -> str:
        """Add a slide to a .pptx. Layout is one of:
        'title', 'title_content', 'section', 'two_content',
        'comparison', 'blank'."""
        try:
            _require_lib("pptx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            prs = Presentation(str(p))
            layout_map = {
                "title": 0,
                "title_content": 1,
                "section": 2,
                "two_content": 3,
                "comparison": 4,
                "title_only": 5,
                "blank": 6,
            }
            idx = layout_map.get((layout or "title").lower(), 1)
            if idx >= len(prs.slide_layouts):
                idx = 1  # safe default
            slide_layout = prs.slide_layouts[idx]
            slide = prs.slides.add_slide(slide_layout)
            # Populate placeholders if they exist on this layout.
            for ph in slide.placeholders:
                ph_idx = ph.placeholder_format.idx
                if ph_idx == 0 and title is not None:
                    ph.text = title
                elif ph_idx == 1 and subtitle is not None:
                    ph.text = subtitle
            prs.save(str(p))
            return f"[ADDED SLIDE] {path} (layout={layout}, slide #{len(prs.slides)})"
        except Exception as e:
            return f"[OFFICE ERROR] add_slide failed: {e}"

    def add_text(self, path: str, slide: int, text: str,
                 x: str = "1in", y: str = "1in",
                 w: str = "8in", h: str = "1in",
                 bold: bool = False, italic: bool = False,
                 color: Optional[str] = None,
                 size: Optional[int] = None,
                 align: Optional[str] = None) -> str:
        """Add a text box to a slide. `slide` is 1-indexed."""
        try:
            _require_lib("pptx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from pptx import Presentation
            from pptx.util import Emu, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            prs = Presentation(str(p))
            if slide < 1 or slide > len(prs.slides):
                return f"[OFFICE ERROR] slide index out of range: {slide}"
            target = prs.slides[slide - 1]
            left = _parse_emu(x) or Emu(914400)
            top = _parse_emu(y) or Emu(914400)
            width = _parse_emu(w) or Emu(914400 * 8)
            height = _parse_emu(h) or Emu(914400)
            txbox = target.shapes.add_textbox(left, top, width, height)
            tf = txbox.text_frame
            tf.text = text
            para = tf.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            run.text = text
            if bold:
                run.font.bold = True
            if italic:
                run.font.italic = True
            if size:
                run.font.size = Pt(int(size))
            if color:
                hexc = _parse_color(color)
                if hexc:
                    run.font.color.rgb = RGBColor.from_string(hexc)
            if align:
                align_map = {
                    "left": PP_ALIGN.LEFT,
                    "center": PP_ALIGN.CENTER,
                    "centre": PP_ALIGN.CENTER,
                    "right": PP_ALIGN.RIGHT,
                    "justify": PP_ALIGN.JUSTIFY,
                }
                para.alignment = align_map.get(align.lower(), PP_ALIGN.LEFT)
            prs.save(str(p))
            return f"[ADDED TEXT] {path} slide {slide}"
        except Exception as e:
            return f"[OFFICE ERROR] add_text failed: {e}"

    def add_shape(self, path: str, slide: int,
                  shape_type: str = "rectangle",
                  x: str = "1in", y: str = "1in",
                  w: str = "2in", h: str = "1in",
                  text: Optional[str] = None,
                  fill_color: Optional[str] = None,
                  line_color: Optional[str] = None) -> str:
        """Add an auto shape to a slide. shape_type: rectangle, rounded,
        ellipse, triangle, arrow, etc."""
        try:
            _require_lib("pptx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from pptx import Presentation
            from pptx.util import Emu, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            prs = Presentation(str(p))
            if slide < 1 or slide > len(prs.slides):
                return f"[OFFICE ERROR] slide index out of range: {slide}"
            target = prs.slides[slide - 1]
            shape_map = {
                "rectangle": MSO_SHAPE.RECTANGLE,
                "rect": MSO_SHAPE.RECTANGLE,
                "rounded": MSO_SHAPE.ROUNDED_RECTANGLE,
                "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
                "ellipse": MSO_SHAPE.OVAL,
                "oval": MSO_SHAPE.OVAL,
                "circle": MSO_SHAPE.OVAL,
                "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
                "arrow": MSO_SHAPE.RIGHT_ARROW,
                "right_arrow": MSO_SHAPE.RIGHT_ARROW,
                "diamond": MSO_SHAPE.DIAMOND,
                "pentagon": MSO_SHAPE.PENTAGON,
                "hexagon": MSO_SHAPE.HEXAGON,
                "star": MSO_SHAPE.STAR_5_POINT,
            }
            mso = shape_map.get((shape_type or "rectangle").lower(),
                                MSO_SHAPE.RECTANGLE)
            left = _parse_emu(x) or Emu(914400)
            top = _parse_emu(y) or Emu(914400)
            width = _parse_emu(w) or Emu(914400 * 2)
            height = _parse_emu(h) or Emu(914400)
            shape = target.shapes.add_shape(mso, left, top, width, height)
            if fill_color:
                hexc = _parse_color(fill_color)
                if hexc:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor.from_string(hexc)
            if line_color:
                hexc = _parse_color(line_color)
                if hexc:
                    shape.line.color.rgb = RGBColor.from_string(hexc)
            if text is not None:
                shape.text_frame.text = text
            prs.save(str(p))
            return f"[ADDED SHAPE] {path} slide {slide} ({shape_type})"
        except Exception as e:
            return f"[OFFICE ERROR] add_shape failed: {e}"

    # ── Cross-format operations ────────────────────────────────────

    def find_replace(self, path: str, find: str, replace: str,
                     sheet: Optional[str] = None,
                     slide: Optional[int] = None) -> str:
        """Find and replace text in any format. Returns the number of
        replacements made.

        For .xlsx, scope by `sheet` (or all sheets if None).
        For .pptx, scope by `slide` (1-indexed) or all slides if None.
        For .docx, scans paragraphs and table cells document-wide.
        """
        try:
            fmt = _detect_format(path)
            _require_lib(fmt)
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            if not find:
                return "[OFFICE ERROR] find is empty"
            if fmt == "docx":
                return self._find_replace_docx(p, find, replace)
            if fmt == "xlsx":
                return self._find_replace_xlsx(p, find, replace, sheet)
            return self._find_replace_pptx(p, find, replace, slide)
        except Exception as e:
            return f"[OFFICE ERROR] find_replace failed: {e}"

    def _find_replace_docx(self, p: Path, find: str, replace: str) -> str:
        from docx import Document
        doc = Document(str(p))
        count = 0
        for para in doc.paragraphs:
            if find in para.text:
                # Rebuild runs: simple approach — set the first run's
                # text to the replaced whole-paragraph text and clear
                # the rest. This loses intra-paragraph formatting but
                # is robust. For finer control users can use str_replace
                # on the underlying XML via run_code.
                new_text = para.text.replace(find, replace)
                count += para.text.count(find)
                if para.runs:
                    para.runs[0].text = new_text
                    for r in para.runs[1:]:
                        r.text = ""
        # Also scan tables.
        for tbl in doc.tables:
            for row in tbl.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        if find in para.text:
                            count += para.text.count(find)
                            new_text = para.text.replace(find, replace)
                            if para.runs:
                                para.runs[0].text = new_text
                                for r in para.runs[1:]:
                                    r.text = ""
        doc.save(str(p))
        return f"[FIND/REPLACE] {p.name}: {count} replacement(s)"

    def _find_replace_xlsx(self, p: Path, find: str, replace: str,
                            sheet: Optional[str]) -> str:
        from openpyxl import load_workbook
        wb = load_workbook(str(p))
        count = 0
        sheets = [wb[sheet]] if sheet else wb.worksheets
        for ws in sheets:
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None and isinstance(cell.value, str) and find in cell.value:
                        count += cell.value.count(find)
                        cell.value = cell.value.replace(find, replace)
        wb.save(str(p))
        wb.close()
        return f"[FIND/REPLACE] {p.name}: {count} replacement(s)"

    def _find_replace_pptx(self, p: Path, find: str, replace: str,
                            slide: Optional[int]) -> str:
        from pptx import Presentation
        prs = Presentation(str(p))
        count = 0
        slides = [prs.slides[slide - 1]] if slide else list(prs.slides)
        for sl in slides:
            for shape in sl.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if find in run.text:
                            count += run.text.count(find)
                            run.text = run.text.replace(find, replace)
        prs.save(str(p))
        return f"[FIND/REPLACE] {p.name}: {count} replacement(s)"

    def save_as(self, path: str, new_path: str) -> str:
        """Save a copy of an office document under a new path. The
        source format and target format must match (this is NOT a
        format converter)."""
        try:
            fmt_src = _detect_format(path)
            fmt_dst = _detect_format(new_path)
            if fmt_src != fmt_dst:
                return (
                    f"[OFFICE ERROR] format mismatch: source is .{fmt_src}, "
                    f"target is .{fmt_dst}. save_as does not convert formats."
                )
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            np = self._resolve(new_path)
            np.parent.mkdir(parents=True, exist_ok=True)
            # Use the appropriate library's save mechanism so the file
            # is well-formed (not just a binary copy).
            if fmt_src == "docx":
                from docx import Document
                Document(str(p)).save(str(np))
            elif fmt_src == "xlsx":
                from openpyxl import load_workbook
                wb = load_workbook(str(p))
                wb.save(str(np))
                wb.close()
            else:
                from pptx import Presentation
                Presentation(str(p)).save(str(np))
            return f"[SAVED AS] {new_path}"
        except Exception as e:
            return f"[OFFICE ERROR] save_as failed: {e}"

    # ── Bulk fill (convenience for agents) ─────────────────────────

    def fill_sheet(self, path: str, sheet: str,
                   data: List[List[Any]],
                   start_cell: str = "A1") -> str:
        """Fill a rectangular block of cells starting at start_cell.
        ``data`` is a list of rows; each row is a list of values.
        Useful for dumping a structured result (e.g. a CSV-like list)
        into a sheet in a single tool call instead of N set_cell calls.
        """
        try:
            _require_lib("xlsx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from openpyxl import load_workbook
            from openpyxl.utils import range_boundaries
            from openpyxl.utils.cell import coordinate_from_string, column_index_from_string
            wb = load_workbook(str(p))
            if sheet not in wb.sheetnames:
                wb.close()
                return f"[OFFICE ERROR] sheet not found: {sheet}"
            ws = wb[sheet]
            col_str, row = coordinate_from_string(start_cell)
            start_col = column_index_from_string(col_str)
            start_row = int(row)
            for i, row_data in enumerate(data):
                for j, val in enumerate(row_data):
                    ws.cell(row=start_row + i,
                            column=start_col + j,
                            value=val)
            wb.save(str(p))
            wb.close()
            return (f"[FILLED SHEET] {path} / {sheet} "
                    f"({len(data)} rows x {max(len(r) for r in data) if data else 0} cols, "
                    f"from {start_cell})")
        except Exception as e:
            return f"[OFFICE ERROR] fill_sheet failed: {e}"

    def fill_table(self, path: str, table_index: int,
                   data: List[List[Any]]) -> str:
        """Replace the contents of an existing table in a .docx.
        ``table_index`` is 1-indexed. Cells beyond the table's current
        dimensions are ignored; missing cells stay empty."""
        try:
            _require_lib("docx")
            p = self._resolve(path)
            if not p.exists():
                return f"[OFFICE ERROR] file not found: {path}"
            from docx import Document
            doc = Document(str(p))
            if table_index < 1 or table_index > len(doc.tables):
                return f"[OFFICE ERROR] table index out of range: {table_index}"
            tbl = doc.tables[table_index - 1]
            rows = len(tbl.rows)
            cols = len(tbl.columns)
            for i, row_data in enumerate(data[:rows]):
                for j, val in enumerate(row_data[:cols]):
                    tbl.cell(i, j).text = str(val) if val is not None else ""
            doc.save(str(p))
            return f"[FILLED TABLE] {path} table#{table_index} ({rows}x{cols})"
        except Exception as e:
            return f"[OFFICE ERROR] fill_table failed: {e}"


# ── Tool schema for the system prompt ──────────────────────────────────

OFFICE_TOOL_SCHEMA = """
# v1.2.0: Office Worker tools (office section only).
# Use these to CREATE, INSPECT, and EDIT .docx / .xlsx / .pptx files
# directly — no need to shell out or generate Python code.
# All file paths are workspace-relative (or absolute inside workspace).
# Every tool returns a short status string like "[CREATED] report.docx".

{"tool": "office_create", "args": {"path": "report.docx", "template": "blank"}}
{"tool": "office_view", "args": {"path": "report.docx", "mode": "outline"}}
  # mode: "outline" (default) | "text" | "stats"

# Word (.docx)
{"tool": "office_add_paragraph", "args": {"path": "report.docx", "text": "Executive summary...", "style": "Normal", "bold": false, "italic": false, "color": null, "size": null}}
  # style: "Normal" | "Heading 1" | "Heading 2" | "List Bullet" | "List Number" | ...
{"tool": "office_add_heading", "args": {"path": "report.docx", "text": "Chapter 1", "level": 1}}
{"tool": "office_add_table", "args": {"path": "report.docx", "rows": 3, "cols": 2, "data": [["Name", "Score"], ["Alice", 95], ["Bob", 87]], "header": true}}
{"tool": "office_fill_table", "args": {"path": "report.docx", "table_index": 1, "data": [["A", "B"], ["1", "2"]]}}

# Excel (.xlsx)
{"tool": "office_add_sheet", "args": {"path": "data.xlsx", "name": "Summary"}}
{"tool": "office_set_cell", "args": {"path": "data.xlsx", "sheet": "Sheet1", "cell": "A1", "value": "Name"}}
  # value may be a string, number, boolean, or a formula like "=SUM(A2:A10)"
{"tool": "office_set_cell_format", "args": {"path": "data.xlsx", "sheet": "Sheet1", "cell": "A1", "bold": true, "italic": false, "font_color": "FF0000", "bg_color": "FFFF00", "font_size": 14, "align": "center"}}
  # colors: "#FF0000" | "FF0000" | "red" | "rgb(255,0,0)"
{"tool": "office_add_chart", "args": {"path": "data.xlsx", "sheet": "Sheet1", "chart_type": "bar", "data_range": "A1:B10", "anchor": "D2", "title": "Revenue"}}
  # chart_type: "bar" | "line" | "pie" | "scatter"
{"tool": "office_fill_sheet", "args": {"path": "data.xlsx", "sheet": "Sheet1", "data": [["Name", "Score"], ["Alice", 95]], "start_cell": "A1"}}

# PowerPoint (.pptx)
{"tool": "office_add_slide", "args": {"path": "deck.pptx", "layout": "title", "title": "Q4 Report", "subtitle": "Revenue grew 25%"}}
  # layout: "title" | "title_content" | "section" | "two_content" | "comparison" | "title_only" | "blank"
{"tool": "office_add_text", "args": {"path": "deck.pptx", "slide": 1, "text": "Key insight: ...", "x": "1in", "y": "2in", "w": "8in", "h": "1in", "bold": true, "size": 24, "color": "FFFFFF", "align": "center"}}
  # dimensions: "1in", "2cm", "12pt", "96px" (parsed into EMU internally)
{"tool": "office_add_shape", "args": {"path": "deck.pptx", "slide": 1, "shape_type": "rectangle", "x": "1in", "y": "3in", "w": "4in", "h": "1in", "text": "Click here", "fill_color": "1A1A2E", "line_color": "FFFFFF"}}
  # shape_type: "rectangle" | "rounded" | "ellipse" | "triangle" | "arrow" | "diamond" | "star" | ...

# Cross-format
{"tool": "office_find_replace", "args": {"path": "report.docx", "find": "TBD", "replace": "Final", "sheet": null, "slide": null}}
  # sheet: optional sheet name (xlsx only); slide: optional 1-indexed slide (pptx only)
{"tool": "office_save_as", "args": {"path": "report.docx", "new_path": "report_v2.docx"}}
  # source and target formats MUST match (no format conversion)
""".strip()


# ── Tool name enum (mirrors ToolName in agent_runtime.py) ──────────────
# Listed here as strings so the parent agent_runtime.py can register
# them without a circular import.

OFFICE_TOOL_NAMES: Tuple[str, ...] = (
    "office_create",
    "office_view",
    "office_add_paragraph",
    "office_add_heading",
    "office_add_table",
    "office_fill_table",
    "office_add_sheet",
    "office_set_cell",
    "office_set_cell_format",
    "office_add_chart",
    "office_fill_sheet",
    "office_add_slide",
    "office_add_text",
    "office_add_shape",
    "office_find_replace",
    "office_save_as",
)


# ── Office-section system prompt suffix ────────────────────────────────

OFFICE_SYSTEM_SUFFIX = """
# Office Worker Mode

You are running in OFFICE WORKER mode. You have access to a dedicated
set of `office_*` tools for creating and editing .docx, .xlsx, and
.pptx files directly — without shelling out, without generating Python
boilerplate, without leaving the agent loop.

## When to use office_* tools
- User asks to create / edit / inspect a Word, Excel, or PowerPoint
  document → use office_* tools. ALWAYS prefer them over write_file +
  asking the user to install python-docx themselves.
- User asks to "make a report", "build a deck", "fill a spreadsheet"
  → use office_create first, then the appropriate add_* / set_* tools.

## Office Worker workflow
1. office_view(path, mode="outline") to inspect an existing file, or
   office_create(path) to start a new one.
2. Build up the document step by step: add_heading, add_paragraph,
   add_table for .docx; add_sheet, set_cell, fill_sheet, add_chart
   for .xlsx; add_slide, add_text, add_shape for .pptx.
3. office_view(path, mode="outline") again to verify the structure
   before reporting back to the user.
4. final_answer should summarise: what file was created/edited, what
   sections/sheets/slides it contains, and any follow-ups (e.g.
   "open in Word to verify pagination").

## Important constraints
- All file paths are workspace-relative. The workspace sandbox applies
  — you cannot write outside it.
- office_save_as does NOT convert formats (.docx → .pdf is not
  supported; for that, suggest the user open in Word/LibreOffice).
- For find/replace across runs in .docx, intra-paragraph formatting
  may be flattened. For surgical edits, use str_replace on the
  document.xml part of the .docx zip — but in practice the
  office_find_replace tool is the right choice for 95% of cases.
- If a library is missing (python-docx / openpyxl / python-pptx), the
  tool returns a clear error telling the user to install it. Do NOT
  try to vendor the library yourself.
- You can still use read_file / write_file / run_code / search_project
  / list_files etc. — the office_* tools are ADDITIVE, not a
  replacement for the general toolkit. Use them when the task is
  specifically about Office documents.
""".strip()


__all__ = [
    "OfficeWorker",
    "OFFICE_TOOL_SCHEMA",
    "OFFICE_TOOL_NAMES",
    "OFFICE_SYSTEM_SUFFIX",
]
